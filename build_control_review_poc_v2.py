#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
File       : create_enhanced_pitch_deck.py
Author     : Gemini
Created    : 2025-08-18
Description: A final, refined script to generate a visually compelling
             PowerPoint presentation for an AI-driven controls monitoring pitch.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR


def create_enhanced_pitch_deck(output_path="AI_Controls_Pitch_Enhanced.pptx"):
    """
    Generates a more visually appealing and persuasive pitch deck.
    """
    prs = Presentation()

    # Define a clean color palette for a professional look
    # Using RGB values for consistency across different themes
    COLOR_DARK_BLUE = RGBColor(38, 50, 56)
    COLOR_ACCENT_BLUE = RGBColor(41, 128, 185)
    COLOR_LIGHT_GRAY = RGBColor(236, 240, 241)

    # === Slide 1: Title Slide (More Dynamic) ===
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Style the title text
    title.text = "AI-Driven Controls Monitoring"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT_BLUE

    # Style the subtitle text for impact
    subtitle.text = (
        "Transforming Risk Management with Intelligent Automation\n\n"
        "Presented by: Ritesh"
    )
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE

    # === Slide 2: Problem Statement (More Visual) ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Challenge: A Reactive & Inefficient Process"

    # Use a list for the body text for better readability
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()  # Clear existing content to rebuild

    p = tf.paragraphs[0]
    p.text = "Manual control reviews are slow, prone to error, and costly."
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_DARK_BLUE

    # Add key pain points as separate, bolded bullets
    for text in [
        "⏰  Time-consuming & labor-intensive",
        "⚖️  Inconsistent assessments lead to audit risk",
        "📈  Lack of scalability with growing data volume",
        "🚨  Reactive, not proactive monitoring"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_BLUE

    # === Slide 3: Business Value Proposition (Clear & Concise) ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Value: Proactive & Automated Assurance"

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "We deliver a comprehensive, automated solution that provides quantifiable ROI."
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_DARK_BLUE

    for text in [
        "⚡️  **Efficiency:** Up to a 75% reduction in manual review effort",
        "🎯  **Accuracy:** Consistent, data-driven assessments",
        "🛡️  **Compliance:** Improved risk posture and faster audit cycles",
        "🚀  **Scalability:** Monitor all controls, all the time"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = False  # Bolding is done inline with the text
        p.font.color.rgb = COLOR_DARK_BLUE

    # === Slide 4: Solution Architecture (More Professional Diagram) ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout for custom content

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Solution Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Recreate the architecture boxes with a more modern look
    box_w, box_h = Inches(2.2), Inches(1)

    # Create the first box with a custom style
    box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), box_w, box_h)
    box1.text = "Data Ingestion"
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    box1.line.fill.background()  # No border
    box1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    box1.text_frame.paragraphs[0].font.size = Pt(16)

    # Add a short description below
    desc1_box = slide.shapes.add_textbox(Inches(1), Inches(3.1), box_w, Inches(0.5))
    desc1_box.text_frame.text = "Pulling data from logs, reports, etc."
    desc1_box.text_frame.paragraphs[0].font.size = Pt(10)

    # Create the second box with a different color/style
    box2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.7), Inches(2), box_w, box_h)
    box2.text = "AI-Driven Analysis"
    box2.fill.solid()
    box2.fill.fore_color.rgb = COLOR_DARK_BLUE
    box2.line.fill.background()
    box2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    box2.text_frame.paragraphs[0].font.size = Pt(16)

    desc2_box = slide.shapes.add_textbox(Inches(3.7), Inches(3.1), box_w, Inches(0.5))
    desc2_box.text_frame.text = "NER, Anomaly Scoring, Prompt Logic"
    desc2_box.text_frame.paragraphs[0].font.size = Pt(10)

    # Create the third box
    box3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(2), box_w, box_h)
    box3.text = "Unified Reporting"
    box3.fill.solid()
    box3.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    box3.line.fill.background()
    box3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    box3.text_frame.paragraphs[0].font.size = Pt(16)

    desc3_box = slide.shapes.add_textbox(Inches(6.4), Inches(3.1), box_w, Inches(0.5))
    desc3_box.text_frame.text = "Real-time dashboards & audit logs"
    desc3_box.text_frame.paragraphs[0].font.size = Pt(10)

    # Add arrows between the new boxes for flow
    # This is a basic arrow, but a more complex graphic would be better
    left, top, w, h = Inches(3.2), Inches(2.4), Inches(0.5), Inches(0.2)
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_DARK_BLUE

    left, top, w, h = Inches(5.9), Inches(2.4), Inches(0.5), Inches(0.2)
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_DARK_BLUE

    # === Slide 5: The End-to-End AI Pipeline ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The End-to-End AI Pipeline"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Our `run_pipeline.py` script orchestrates the entire process, from data to insights."
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_DARK_BLUE

    for text in [
        "1. **Generate Inputs:** Creates a synthetic dataset for training and evaluation.",
        "2. **Train Models:** Trains specialized NER and sentiment models for your use case.",
        "3. **Run Inference:** Applies the trained models to new data to extract entities and sentiment.",
        "4. **Composite Model:** Combines the outputs of multiple models to create a holistic view.",
        "5. **Apply Business Logic:** Uses the model outputs to trigger specific, predefined business rules."
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = False
        p.font.color.rgb = COLOR_DARK_BLUE

    # === Slide 6: Results & Metrics (More engaging chart) ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quantifiable Results"

    chart_data = ChartData()
    chart_data.categories = ['Control Accuracy', 'Effort Reduction']
    chart_data.add_series('PoC Metrics', (90, 75))

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = chart_frame.chart

    # Style the chart for better visual appeal
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = "PoC Metrics: A Snapshot"

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False  # Remove gridlines for a cleaner look
    value_axis.maximum_scale = 100
    value_axis.minimum_scale = 0

    # Add a data label to the columns
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(12)
    data_labels.font.color.rgb = COLOR_DARK_BLUE

    # === Slide 7: Next Steps ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Path Forward"

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    # Create an ordered list for a clear action plan
    for text in [
        "1. Pilot a use case with a specific team or department.",
        "2. Integrate with your existing GRC and audit systems.",
        "3. Deploy a custom, real-time dashboard for continuous monitoring.",
        "4. Scale the solution across the entire enterprise."
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK_BLUE

    # === Slide 8: Thank You ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank You"
    slide.shapes.placeholders[1].text = (
        "Let's work together to redefine control assurance.\n\n"
        "Contact: Ritesh\n"
        "Email: your.email@example.com"
    )

    prs.save(output_path)
    print(f"✅ Enhanced presentation generated: {output_path}")

if __name__ == "__main__":
    create_enhanced_pitch_deck()
