#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
File       : create_enhanced_pitch_deck.py
Author     : Copilot
Created    : 2025-08-18
Description: Generates a polished PowerPoint pitch deck for AI-driven
             controls monitoring, with global stats, images, and references.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- CONFIG ---
OUTPUT_PATH = "AI_Controls_Pitch_Enhanced_V2.pptx"
IMG_DIR      = "images"
STATS_IMG    = os.path.join(IMG_DIR, "market_growth.png")
PIPE_IMG     = os.path.join(IMG_DIR, "pipeline_diagram.png")

# --- STYLE HELPERS ---
def style_title(shape):
    """Apply title font styling to a text box or placeholder."""
    for p in shape.text_frame.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(38, 50, 56)
        p.alignment = PP_ALIGN.LEFT

def style_body(shape, size=18):
    """Apply body font styling to a text box or placeholder."""
    for p in shape.text_frame.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(38, 50, 56)
        p.level = 0

def set_solid_fill(shape, color):
    """Fill shape with a solid color."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def clear_line(shape):
    """Remove shape outline."""
    shape.line.fill.background()

# --- MAIN GENERATOR ---
def create_enhanced_pitch_deck(output_path=OUTPUT_PATH):
    prs = Presentation()
    ACCENT = RGBColor(41, 128, 185)
    DARK   = RGBColor(38, 50, 56)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI-Driven Controls Monitoring"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(48)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT

    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Transforming Risk Management with Intelligent Automation\n\n"
        "Presented by: Ritesh"
    )
    st = subtitle.text_frame.paragraphs[0]
    st.font.size = Pt(28)
    st.font.color.rgb = DARK

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Challenge: Reactive & Inefficient"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.add_paragraph().text = "Manual control reviews are slow, error-prone, and costly."
    for bullet in [
        "⏰  Time-consuming & labour-intensive",
        "⚖️  Inconsistent assessments → audit risk",
        "📈  No real-time visibility or scalability",
        "🚨  Reactive rather than proactive monitoring"
    ]:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK

    # Slide 3: Business Value Proposition
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Value: Proactive Assurance"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.add_paragraph().text = "A unified ML-powered pipeline delivering quantifiable ROI."
    for bullet in [
        "⚡️  Up to 75% reduction in manual effort",
        "🎯  Consistent, data-driven control assessments",
        "🛡️  Faster audit cycles & improved compliance",
        "🚀  Real-time, scalable monitoring of all controls"
    ]:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = DARK

    # Slide 4: Solution Architecture (fixed fill usage)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
    )
    title_box.text_frame.text = "Solution Architecture"
    style_title(title_box)

    steps = ["Data Ingestion", "AI-Driven Analysis", "Unified Reporting"]
    box_w, box_h = Inches(2.2), Inches(1)
    for i, label in enumerate(steps):
        x = Inches(1 + i * 2.8)
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(1.8), box_w, box_h
        )
        # apply solid fill & clear outline
        color = ACCENT if (i % 2 == 0) else DARK
        set_solid_fill(box, color)
        clear_line(box)

        tf = box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(16)

    # Arrows between the boxes
    for i in range(len(steps) - 1):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(1 + i * 2.8) + box_w, Inches(2.3),
            Inches(0.6), Inches(0.4)
        )
        set_solid_fill(arrow, DARK)
        clear_line(arrow)

    # Slide 5: End-to-End AI Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "End-to-End AI Pipeline"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for step in [
        "1. Generate inputs from logs, reports",
        "2. Train NER & sentiment models",
        "3. Run inference on live data",
        "4. Composite model for PASS/FAIL flags",
        "5. Unified audit report & dashboards"
    ]:
        p = body.add_paragraph()
        p.text = step
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK

    # Slide 6: Global Market Stats
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
    )
    title_box.text_frame.text = "Global AI Monitoring Market"
    style_title(title_box)

    stats_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1), Inches(5), Inches(3)
    )
    tf = stats_box.text_frame
    tf.text = "• Market projected to reach USD 12.5 B by 2027 (CAGR 26.2%)"
    for bullet in [
        "• 78 % of orgs use AI in ≥ 1 function (2025)",
        "• Real-time AI monitoring reduces outages by 40 %",
        "• Drift detection cuts model retrains by 30 % annually"
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK

    if os.path.exists(STATS_IMG):
        slide.shapes.add_picture(
            STATS_IMG, Inches(6), Inches(1), width=Inches(3)
        )

    # Slide 7: Sample Pipeline Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
    )
    title_box.text_frame.text = "Pipeline in Action"
    style_title(title_box)

    if os.path.exists(PIPE_IMG):
        slide.shapes.add_picture(
            PIPE_IMG, Inches(1), Inches(1), width=Inches(8)
        )

    # Slide 8: Quantifiable Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quantifiable PoC Metrics"
    data = ChartData()
    data.categories = ["Control Accuracy", "Effort Saved"]
    data.add_series("Metrics", (90, 75))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(8), Inches(4),
        data
    ).chart
    chart.value_axis.maximum_scale = 100
    chart.has_legend = False
    chart.plots[0].has_data_labels = True

    # Slide 9: Study References
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Study References"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for ref in [
        "Urrea C. et al., “AI-Driven & Bio-Inspired Control,” Machines, 2025",
        "Zohuri B., “AI/ML Adaptive Control Applications,” JMSET, 2024",
        "Komala M.H. et al., “AI in Power Electronics,” IJCRT, 2022"
    ]:
        p = body.add_paragraph()
        p.text = ref
        p.font.size = Pt(16)
        p.font.color.rgb = DARK

    # Slide 10: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Path Forward"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for step in [
        "1. Pilot with a control team & integrate GRC tools",
        "2. Fine-tune thresholds via config.yaml",
        "3. Deploy real-time dashboards & drift alerts",
        "4. Scale across the enterprise"
    ]:
        p = body.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = DARK

    # Slide 11: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank You"
    slide.shapes.placeholders[1].text = (
        "Let's redefine control assurance with AI-driven automation.\n\n"
        "Contact: Ritesh\nemail@yourdomain.com"
    )

    prs.save(output_path)
    print(f"✅ Presentation successfully generated: {output_path}")


if __name__ == "__main__":
    create_enhanced_pitch_deck()
