#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
File       : build_control_review_poc.py
Author     : Ritesh
Created    : 2025-08-18
Description: <SHORT DESCRIPTION HERE>
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

def create_control_review_poc(output_path="Automated_Control_Review_POC.pptx"):
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automated Control Effectiveness Review"
    slide.placeholders[1].text = (
        "Config-Driven ML Pipeline to Optimize Controls\n"
        "and Reduce Workforce Overhead\n\n"
        "Presented by: Ritesh"
    )

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Manual control reviews are time-consuming and error-prone"
    for bullet in [
        "Inconsistent assessments lead to audit risk",
        "High labor costs and resource bottlenecks",
        "Lack of real-time visibility into control effectiveness"
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 3: Business Value Proposition
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Value Proposition"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Automated, consistent control assessments"
    for bullet in [
        "Up to 75% reduction in manual review effort",
        "Faster audit cycles with real-time reports",
        "Improved compliance through standardized validation"
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 4: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Solution Overview"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Modular pipeline for control data ingestion and validation"
    for bullet in [
        "Control entity extraction via NER modules",
        "Risk and anomaly scoring with ML models",
        "Prompt-aware business logic for PASS/FAIL flags",
        "Unified audit summary: counts, accuracy, mismatches"
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 5: Solution Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Solution Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)

    # Architecture boxes
    left, top, w, h = Inches(1), Inches(1.2), Inches(2.2), Inches(1)
    labels = ["Control Entity Extraction", "Risk & Anomaly Scoring", "Control Effectiveness Validation"]
    boxes = []
    for idx, lbl in enumerate(labels):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + idx * Inches(3), top, w, h
        )
        box.text = lbl
        boxes.append(box)

    # Arrows between boxes
    for i in range(len(boxes) - 1):
        src = boxes[i]
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            src.left + src.width, src.top + h / 3,
            Inches(1), Inches(0.5)
        )
        # fill arrow solid black
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)
        # set arrow border to black
        arrow.line.color.rgb = RGBColor(0, 0, 0)

    # Slide 6: Sample Use Case
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sample Use Case"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Input: Control log entries with metadata"
    for bullet in [
        "Process: NER → Scoring → Prompt-based rule engine",
        "Output: PASS/FAIL control flags & reasoning",
        "Artifacts: CSV output, human-readable summary, audit log"
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 7: Results & Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Results & Metrics"
    chart_data = ChartData()
    chart_data.categories = ['Control Review Accuracy', 'Review Effort Saved']
    chart_data.add_series('Metrics', (90, 75))

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.value_axis.maximum_scale = 100
    chart.has_legend = False

    # Slide 8: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Next Steps"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Integrate with enterprise audit systems"
    for bullet in [
        "Fine-tune thresholds via config.yaml",
        "Deploy rollback triggers for minor/major drifts",
        "Build real-time dashboard for controls monitoring"
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 9: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = (
        "Let’s redefine control assurance with AI-driven automation.\n\n"
        "Contact: ritesh@yourdomain.com"
    )

    # Save the deck
    prs.save(output_path)
    print(f"✅ Presentation generated: {output_path}")

if __name__ == "__main__":
    create_control_review_poc()
